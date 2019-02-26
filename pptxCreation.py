####################################################################################################
# Author: Edgar Oregel
# Date: 6/20/18
# File: interface.py
# Description: Controls the UI of the application and its components
# References: https://python-pptx.readthedocs.io/en/latest/user/presentations.html 
####################################################################################################


#########
#Imports#
#########
import os
import re
import logging
from pptx import Presentation


##################
#Global Variables#
##################

############
#Main Class#
############
class PPTX():
    def __init__(self):
        #Instantiate logger
        self.logger = logging.getLogger()
        handler = logging.StreamHandler()
        formatter = logging.Formatter(
        '%(asctime)s %(name)-12s %(levelname)-8s %(message)s')
        handler.setFormatter(formatter)
        self.logger.addHandler(handler)
        self.logger.setLevel(logging.DEBUG)


    def createPresentation(self, name, slides=None):
        #check if powerpoint name exists
        return
    def getPowerPointSlides(self, name):
        #open powerpoint
        #get slides
        #return them 
        return
    
    def openPresentation(self, name):
        try:
            _file = open(name, 'w')
        except:
            self.logger.error("[openPresentation] Couldn't open presentation: {}".format(name))
            sys.exit(1)
        pptx = Presentation(_file)


###################
#Utility Functions#
###################
